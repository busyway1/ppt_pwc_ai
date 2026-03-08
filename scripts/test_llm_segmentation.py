#!/usr/bin/env python3
"""
LLM 세그멘테이션 A/B 비교 테스트.

규칙 기반 결과와 LLM 보강 결과를 비교하여 품질 지표를 측정한다.

Usage:
    python scripts/test_llm_segmentation.py input.pptx [--slides 0,1,2]

지표:
  - item_count=1 비율 (목표: <30%, 현재 baseline: 95.2%)
  - body-only zone 비율 (목표: <50%, 현재 baseline: ~90%)
  - card-grid 기본값 비율 (목표: <25%, 현재 baseline: 47.6%)
  - 컴포넌트 타입 다양성
  - relationships 추출률 (목표: >70%)
  - design_intent 존재 비율 (목표: 100%)
  - text_label 다양성 (목표: >60%)
"""

import argparse
import json
import sys
from pathlib import Path

PROJECT_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(PROJECT_ROOT / "tools"))
sys.path.insert(0, str(PROJECT_ROOT / "scripts"))
sys.path.insert(0, str(PROJECT_ROOT / "pipeline"))


def run_comparison(pptx_path, slide_indices=None):
    """규칙 기반 vs LLM 보강 A/B 비교."""
    from inventory import extract_text_inventory, ShapeData
    from extract_design_tokens import extract_design_tokens
    from segment_slide import segment_slide
    from classify_component import classify_component
    from llm_segmentation import enhance_slide_segmentation, merge_llm_with_rules
    from llm_client import is_configured
    from pptx import Presentation

    pptx_path = Path(pptx_path)
    config_path = PROJECT_ROOT / "config" / "component-types.json"
    with open(config_path, "r", encoding="utf-8") as f:
        component_types = json.load(f)["types"]

    tokens = extract_design_tokens(pptx_path)
    prs = Presentation(str(pptx_path))
    slide_w = round(ShapeData.emu_to_inches(prs.slide_width), 2)
    slide_h = round(ShapeData.emu_to_inches(prs.slide_height), 2)
    inventory = extract_text_inventory(pptx_path, prs)

    use_llm = is_configured()
    if not use_llm:
        print("WARNING: LLM not configured. Only rule-based results will be shown.")
        print("Set GENAI_BASE_URL and PwC_LLM_API_KEY in .env to enable LLM.\n")

    # 지표 수집
    rule_metrics = _init_metrics()
    llm_metrics = _init_metrics()

    for slide_id, shapes_dict in inventory.items():
        slide_idx = int(slide_id.split("-")[1])
        if slide_indices and slide_idx not in slide_indices:
            continue

        shapes_list = list(shapes_dict.values())
        print(f"\n{'=' * 60}")
        print(f"Slide {slide_idx} ({len(shapes_list)} shapes)")
        print(f"{'=' * 60}")

        # --- Rule-based ---
        rule_segments = segment_slide(shapes_list, slide_w, slide_h)
        print(f"\n[RULES] {len(rule_segments)} segments:")
        for i, seg in enumerate(rule_segments):
            comp_type = classify_component(seg, component_types, slide_w, slide_h)
            n = len(seg["shapes"])
            _collect_metrics(rule_metrics, seg, comp_type, n)
            print(f"  [{i}] {comp_type} (zone={seg['zone']}, shapes={n})")

        # --- LLM-enhanced ---
        if use_llm:
            try:
                llm_result = enhance_slide_segmentation(
                    shapes_list, slide_w, slide_h, component_types
                )
                llm_segments = merge_llm_with_rules(
                    rule_segments, llm_result, shapes_list, slide_w, slide_h
                )
                print(f"\n[LLM] {len(llm_segments)} segments:")
                for i, seg in enumerate(llm_segments):
                    comp_type = classify_component(
                        seg, component_types, slide_w, slide_h
                    )
                    n = len(seg["shapes"])
                    _collect_metrics(
                        llm_metrics,
                        seg,
                        comp_type,
                        n,
                        has_intent=bool(seg.get("llm_design_intent")),
                        has_tags=bool(seg.get("llm_search_tags")),
                        items=seg.get("llm_items", []),
                    )
                    intent_str = ", ".join(seg.get("llm_design_intent", []))
                    print(
                        f"  [{i}] {comp_type} "
                        f"(zone={seg['zone']}, shapes={n}, "
                        f"items={seg.get('llm_item_count', '?')}, "
                        f"intent=[{intent_str}])"
                    )

                # 관계성
                rels = llm_result.get("relationships", [])
                if rels:
                    print(f"\n  Relationships ({len(rels)}):")
                    for r in rels:
                        print(
                            f"    {r['from_component']}→{r['to_component']} ({r['type']}): {r.get('description', '')}"
                        )
                    llm_metrics["relationships_found"] += len(rels)

                llm_metrics["slides_with_relationships"] += 1 if rels else 0
                print(f"  Slide intent: {llm_result.get('slide_intent', '?')}")

            except Exception as e:
                print(f"\n[LLM] ERROR: {e}")

    # --- 결과 비교 ---
    print(f"\n{'=' * 60}")
    print("A/B Comparison Summary")
    print(f"{'=' * 60}")
    _print_metrics("RULES (Baseline)", rule_metrics)
    if use_llm:
        _print_metrics("LLM Enhanced", llm_metrics)


def _init_metrics():
    return {
        "total_components": 0,
        "item_count_1": 0,
        "body_only_zones": 0,
        "total_zones": 0,
        "card_grid_default": 0,
        "type_distribution": {},
        "relationships_found": 0,
        "slides_with_relationships": 0,
        "has_design_intent": 0,
        "has_search_tags": 0,
        "text_label_diverse": 0,  # hook/data/cta 등 body 외 라벨 있음
        "total_slides": 0,
    }


def _collect_metrics(
    metrics, seg, comp_type, n_shapes, has_intent=False, has_tags=False, items=None
):
    metrics["total_components"] += 1
    if n_shapes <= 1 and seg.get("zone") not in ("decoration", "footer"):
        metrics["item_count_1"] += 1
    if comp_type == "card-grid":
        metrics["card_grid_default"] += 1
    metrics["type_distribution"][comp_type] = (
        metrics["type_distribution"].get(comp_type, 0) + 1
    )
    if has_intent:
        metrics["has_design_intent"] += 1
    if has_tags:
        metrics["has_search_tags"] += 1

    # text_label 다양성 확인
    if items:
        for item in items:
            for zone in item.get("zones", []):
                metrics["total_zones"] += 1
                label = zone.get("text_label", "body")
                if label == "body":
                    metrics["body_only_zones"] += 1
                if label in ("hook", "data", "cta", "label", "caption"):
                    metrics["text_label_diverse"] += 1


def _print_metrics(label, m):
    total = m["total_components"] or 1
    total_zones = m["total_zones"] or 1
    print(f"\n  --- {label} ---")
    print(f"  Total components: {m['total_components']}")
    print(
        f"  item_count=1 ratio: {m['item_count_1']}/{total} ({m['item_count_1'] / total * 100:.1f}%)"
    )
    if m["total_zones"]:
        body_ratio = m["body_only_zones"] / total_zones * 100
        diverse_ratio = m["text_label_diverse"] / total_zones * 100
        print(
            f"  body-only zones: {m['body_only_zones']}/{total_zones} ({body_ratio:.1f}%)"
        )
        print(
            f"  text_label diversity: {m['text_label_diverse']}/{total_zones} ({diverse_ratio:.1f}%)"
        )
    print(
        f"  card-grid default: {m['card_grid_default']}/{total} ({m['card_grid_default'] / total * 100:.1f}%)"
    )
    print(
        f"  design_intent present: {m['has_design_intent']}/{total} ({m['has_design_intent'] / total * 100:.1f}%)"
    )
    print(f"  relationships found: {m['relationships_found']}")
    print(
        f"  Type distribution: {json.dumps(m['type_distribution'], ensure_ascii=False)}"
    )


def main():
    parser = argparse.ArgumentParser(description="LLM segmentation A/B comparison test")
    parser.add_argument("input", help="Input PowerPoint file (.pptx)")
    parser.add_argument(
        "--slides",
        default=None,
        help="Comma-separated slide indices (0-based, default: all)",
    )
    args = parser.parse_args()

    pptx_path = Path(args.input)
    if not pptx_path.exists():
        print(f"Error: File not found: {pptx_path}")
        sys.exit(1)

    slide_indices = None
    if args.slides:
        slide_indices = [int(s.strip()) for s in args.slides.split(",")]

    run_comparison(pptx_path, slide_indices)


if __name__ == "__main__":
    main()
