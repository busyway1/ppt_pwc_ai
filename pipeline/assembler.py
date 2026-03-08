#!/usr/bin/env python3
"""
여러 개별 슬라이드 PPTX 파일을 하나의 프레젠테이션으로 합친다.

Usage:
    python pipeline/assembler.py --slides slide1.pptx slide2.pptx --output final.pptx
    python pipeline/assembler.py --slides slide1.pptx slide2.pptx --output final.pptx --template base.pptx
"""

import argparse
import sys
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu


def _add_image_rel(target_slide, old_rel):
    """이미지/미디어 관계를 대상 슬라이드에 안전하게 추가한다."""
    # python-pptx는 get_or_add가 없으므로 기존 관계를 확인 후 추가
    for existing_rel in target_slide.part.rels.values():
        if (
            existing_rel._target is old_rel._target
            and existing_rel.reltype == old_rel.reltype
        ):
            return existing_rel.rId
    # 새 관계 추가
    return target_slide.part.rels._add(old_rel.reltype, old_rel._target)


def _copy_shapes_to_slide(source_slide, target_slide):
    """원본 슬라이드의 모든 shape를 대상 슬라이드로 복사한다."""
    # 이미지/미디어 관계 수집
    image_rels = {}
    for rel_id, rel in source_slide.part.rels.items():
        if "image" in rel.reltype or "media" in rel.reltype:
            image_rels[rel_id] = rel

    sp_tree = target_slide.shapes._spTree

    for shape in source_slide.shapes:
        el = shape.element
        new_el = deepcopy(el)
        # p:extLst가 없을 수 있으므로 안전하게 append
        ext_lst = sp_tree.findall(
            "{http://schemas.openxmlformats.org/presentationml/2006/main}extLst"
        )
        if ext_lst:
            ext_lst[0].addprevious(new_el)
        else:
            sp_tree.append(new_el)

        # 이미지 shape의 blip 참조 업데이트
        blips = new_el.xpath(".//a:blip[@r:embed]")
        for blip in blips:
            old_rId = blip.get(
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
            )
            if old_rId in image_rels:
                new_rId = _add_image_rel(target_slide, image_rels[old_rId])
                blip.set(
                    "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed",
                    new_rId,
                )

    # 추가 이미지/미디어 관계 복사
    for rel_id, rel in image_rels.items():
        try:
            _add_image_rel(target_slide, rel)
        except Exception:
            pass


def _copy_slide_background(source_slide, target_slide):
    """원본 슬라이드의 배경을 대상 슬라이드로 복사한다."""
    try:
        bg = source_slide.background
        if bg is None or bg.element is None:
            return
        bg_elem = bg.element
        target_bg = deepcopy(bg_elem)
        # 배경 요소에서 이미지 참조가 있으면 관계도 복사
        blips = target_bg.xpath(".//a:blip[@r:embed]")
        for blip in blips:
            old_rId = blip.get(
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
            )
            if old_rId:
                for rel_id, rel in source_slide.part.rels.items():
                    if rel_id == old_rId:
                        new_rId = _add_image_rel(target_slide, rel)
                        blip.set(
                            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed",
                            new_rId,
                        )
                        break
        # 대상 슬라이드에 배경 적용
        target_slide.element.insert(0, target_bg)
    except Exception:
        pass  # 배경 복사 실패 시 무시 (기본 배경 유지)


def assemble_presentation(
    slide_pptx_paths: list[str | Path],
    output_path: str | Path,
    template_path: str | Path | None = None,
) -> Path:
    """여러 PPTX 파일을 하나로 합친다.

    Args:
        slide_pptx_paths: 개별 슬라이드 PPTX 파일 경로 목록.
        output_path: 출력 PPTX 파일 경로.
        template_path: 기반 템플릿 PPTX 파일 경로 (테마/마스터 슬라이드 상속).

    Returns:
        생성된 PPTX 파일의 Path.
    """
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    # 기반 프레젠테이션 생성/열기
    if template_path:
        template_path = Path(template_path)
        if not template_path.exists():
            raise FileNotFoundError(f"Template not found: {template_path}")
        base_prs = Presentation(str(template_path))
        # 템플릿의 기존 슬라이드 제거 (레이아웃/마스터만 유지)
        while len(base_prs.slides) > 0:
            rId = base_prs.slides._sldIdLst[0].rId
            base_prs.part.drop_rel(rId)
            del base_prs.slides._sldIdLst[0]
    else:
        base_prs = Presentation()
        base_prs.slide_width = Emu(int(13.33 * 914400))
        base_prs.slide_height = Emu(int(7.5 * 914400))

    # 기본 레이아웃 (빈 슬라이드)
    blank_layout = None
    for layout in base_prs.slide_layouts:
        if "blank" in layout.name.lower() or layout.name == "빈 화면":
            blank_layout = layout
            break
    if blank_layout is None:
        # 가장 마지막 레이아웃 사용 (일반적으로 blank)
        blank_layout = base_prs.slide_layouts[-1]

    for slide_path in slide_pptx_paths:
        slide_path = Path(slide_path)
        if not slide_path.exists():
            raise FileNotFoundError(f"Slide PPTX not found: {slide_path}")

        src_prs = Presentation(str(slide_path))
        for src_slide in src_prs.slides:
            # 새 빈 슬라이드 추가
            new_slide = base_prs.slides.add_slide(blank_layout)
            # 기본 placeholder 제거
            for shape in list(new_slide.shapes):
                sp = shape.element
                sp.getparent().remove(sp)

            # 배경 + shape 복사
            _copy_slide_background(src_slide, new_slide)
            _copy_shapes_to_slide(src_slide, new_slide)

    base_prs.save(str(output_path))
    return output_path


def main():
    parser = argparse.ArgumentParser(
        description="여러 슬라이드 PPTX를 하나의 프레젠테이션으로 합친다.",
    )
    parser.add_argument(
        "--slides", nargs="+", required=True, help="개별 슬라이드 PPTX 파일 경로들"
    )
    parser.add_argument("--output", required=True, help="출력 PPTX 파일 경로")
    parser.add_argument(
        "--template", default=None, help="기반 템플릿 PPTX (테마/마스터 상속)"
    )

    args = parser.parse_args()

    try:
        output = assemble_presentation(args.slides, args.output, args.template)
        print(f"Assembled {len(args.slides)} slide(s) into: {output}")
    except FileNotFoundError as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
